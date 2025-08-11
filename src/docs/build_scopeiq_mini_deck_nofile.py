from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
SCOPEIQ_BLUE = RGBColor(37, 99, 235)
DARK_GRAY    = RGBColor(35, 47, 62)
LIGHT_GRAY   = RGBColor(248, 249, 249)
AWS_STORAGE  = RGBColor(255, 153, 0)
AWS_DB       = RGBColor(140, 79, 255)
AWS_COMPUTE  = RGBColor(255, 181, 71)
AWS_NET      = RGBColor(30, 115, 190)
AWS_SECURITY = RGBColor(0, 125, 188)
AWS_AI       = RGBColor(254, 203, 47)
AWS_MONITOR  = RGBColor(0, 153, 153)
AWS_GRAY     = RGBColor(95, 107, 109)

def add_watermark(slide):
    tx = slide.shapes.add_textbox(Inches(7.0), Inches(6.7), Inches(3.0), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "ScopeIQ"
    r.font.size = Pt(16)
    r.font.color.rgb = SCOPEIQ_BLUE
    p.alignment = PP_ALIGN.RIGHT

def add_tiny_ai_icon(slide):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.25), Inches(6.8), Inches(0.35), Inches(0.35))
    circ.fill.solid(); circ.fill.fore_color.rgb = SCOPEIQ_BLUE; circ.fill.transparency = 0.2
    circ.line.color.rgb = SCOPEIQ_BLUE
    tx = slide.shapes.add_textbox(Inches(0.25), Inches(6.8), Inches(0.35), Inches(0.35))
    tf = tx.text_frame; tf.text = "AI"
    tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = DARK_GRAY; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_service_box(slide, label, x, y, w, h, color, opacity=0.15, text_color=DARK_GRAY, bold=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.fill.transparency = opacity
    rect.line.color.rgb = color
    tx = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(h-0.2))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = label; r.font.size = Pt(12); r.font.bold = bold; r.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return rect

def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = AWS_GRAY; line.line.width = Pt(1.5)
    return line

prs = Presentation()

# Slide 1: Exec summary
s1 = prs.slides.add_slide(prs.slide_layouts[1])
s1.shapes.title.text = "ScopeIQ Architecture – Executive Summary"
body = s1.shapes.placeholders[1].text_frame; body.clear()
bullets = [
    "Secure multi-tenant SaaS built on AWS serverless architecture",
    "Authentication & RBAC via Amazon Cognito; presigned S3 URLs for secure file access",
    "Scalable backend with AWS Lambda, DynamoDB, and CloudFront CDN",
    "AI-powered processing: Textract, Comprehend, GPT-4, PDF.js",
    "Semantic search via Pinecone vector DB, exposed with AppSync"
]
for i,b in enumerate(bullets):
    p = body.add_paragraph() if i else body.paragraphs[0]
    p.text = b; p.font.size = Pt(14); p.font.color.rgb = DARK_GRAY
add_watermark(s1); add_tiny_ai_icon(s1)

# Slide 2: High-Level Architecture (shapes)
s2 = prs.slides.add_slide(prs.slide_layouts[5])
s2.shapes.title.text = "High-Level Architecture"
# Row: Client
add_service_box(s2, "Client Layer\nReact Web App • PWA • Query Scoping", 0.5, 1.2, 9.0, 0.9, LIGHT_GRAY, opacity=1.0)
# Core
add_service_box(s2, "Amplify / CloudFront", 0.5, 2.4, 2.7, 1.0, AWS_NET)
add_service_box(s2, "Cognito", 3.35, 2.4, 2.2, 1.0, AWS_SECURITY)
add_service_box(s2, "AppSync", 5.8, 2.4, 2.2, 1.0, SCOPEIQ_BLUE, opacity=0.25, text_color=RGBColor(255,255,255))
add_service_box(s2, "Lambda", 8.2, 2.4, 1.3, 1.0, AWS_COMPUTE)
# Data
add_service_box(s2, "Amazon S3\n(Documents & Thumbnails)", 0.5, 3.7, 3.2, 1.0, AWS_STORAGE)
add_service_box(s2, "DynamoDB\n(Multi‑Tenant DB)", 3.9, 3.7, 3.0, 1.0, AWS_DB)
add_service_box(s2, "AI & Search\nOpenAI • Pinecone • PDF.js", 7.1, 3.7, 2.4, 1.0, AWS_AI)
# Monitoring
add_service_box(s2, "CloudWatch\n(Monitoring & Logs)", 0.5, 4.9, 2.7, 0.9, AWS_MONITOR)
# Connectors
add_connector(s2, 1.9, 2.4, 1.9, 1.2)
add_connector(s2, 4.45, 2.4, 4.45, 1.2)
add_connector(s2, 6.9, 2.4, 6.9, 1.2)
add_connector(s2, 6.9, 3.4, 6.9, 2.4)
add_connector(s2, 2.1, 3.7, 2.1, 3.4)
add_connector(s2, 5.4, 3.7, 5.4, 3.4)
add_connector(s2, 8.3, 3.4, 8.3, 2.4)
add_connector(s2, 8.3, 4.7, 8.3, 4.3)
# Summary
tx2 = s2.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9.0), Inches(0.7))
tf2 = tx2.text_frame; tf2.text = ("Amplify/CloudFront deliver the app; Cognito secures access; AppSync provides APIs; "
    "Lambda handles processing; S3/DynamoDB store content; AI/Search (OpenAI, Pinecone, PDF.js) powers insights; CloudWatch monitors.")
tf2.paragraphs[0].font.size = Pt(12); tf2.paragraphs[0].font.color.rgb = DARK_GRAY
add_watermark(s2); add_tiny_ai_icon(s2)

# Slide 3: AI & Data Flow (shapes)
s3 = prs.slides.add_slide(prs.slide_layouts[5])
s3.shapes.title.text = "AI & Data Flow"
def box(text, x, y, w, h, col, op=0.15, white=False):
    return add_service_box(s3, text, x, y, w, h, col, opacity=op, text_color=RGBColor(255,255,255) if white else DARK_GRAY)
box("1) Upload → S3", 0.5, 1.6, 2.5, 0.9, AWS_STORAGE)
box("2) Metadata → DynamoDB", 3.2, 1.6, 2.5, 0.9, AWS_DB)
box("3) Processing → Lambda", 5.9, 1.6, 2.5, 0.9, AWS_COMPUTE)
add_service_box(s3, "4) AI Analysis → Textract • Comprehend • GPT‑4 • PDF.js", 0.5, 2.8, 8.0, 0.9, AWS_AI, opacity=0.10)
box("5) Embedding/Index → Pinecone", 0.5, 4.0, 2.9, 0.9, SCOPEIQ_BLUE, op=0.25, white=True)
box("6) Search API → AppSync", 3.6, 4.0, 2.9, 0.9, SCOPEIQ_BLUE, op=0.25, white=True)
box("7) Monitoring → CloudWatch", 6.7, 4.0, 2.0, 0.9, AWS_MONITOR)
add_connector(s3, 1.75, 2.5, 4.45, 2.5)
add_connector(s3, 4.45, 2.5, 7.15, 2.5)
add_connector(s3, 7.15, 2.5, 7.15, 2.8)
add_connector(s3, 2.0, 4.9, 4.0, 4.9)
tx3 = s3.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9.0), Inches(0.7))
tf3 = tx3.text_frame; tf3.text = ("Files in S3; metadata in DynamoDB; Lambda orchestrates analysis via Textract, Comprehend, GPT‑4, and PDF.js; "
                                  "embeddings land in Pinecone; AppSync exposes semantic search; CloudWatch monitors across the stack.")
tf3.paragraphs[0].font.size = Pt(12); tf3.paragraphs[0].font.color.rgb = DARK_GRAY
add_watermark(s3); add_tiny_ai_icon(s3)

# Slide 4: Simplified Technical Architecture (bullets)
s4 = prs.slides.add_slide(prs.slide_layouts[1])
s4.shapes.title.text = "Simplified Technical Architecture"
body4 = s4.shapes.placeholders[1].text_frame; body4.clear()
for i, pt in enumerate([
    "Client Layer (React Web App, PWA, Query Scoping)",
    "Hosting/CDN: AWS Amplify, CloudFront",
    "Authentication: Amazon Cognito, JWT, RBAC",
    "API Layer: AWS AppSync",
    "Processing: AWS Lambda",
    "Storage: Amazon S3",
    "Database: Amazon DynamoDB",
    "AI & Search: GPT‑4, Textract, Comprehend, Pinecone, PDF.js",
    "Monitoring: Amazon CloudWatch"
]):
    p = body4.add_paragraph() if i else body4.paragraphs[0]
    p.text = pt; p.font.size = Pt(14); p.font.color.rgb = DARK_GRAY
add_watermark(s4); add_tiny_ai_icon(s4)

prs.save("scopeiq_architecture_mini_deck.pptx")
print("Saved scopeiq_architecture_mini_deck.pptx")
