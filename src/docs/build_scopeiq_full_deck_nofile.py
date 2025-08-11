# build_scopeiq_full_deck_nofile.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import os

# ============ Brand & Palette ============
SCOPEIQ_BLUE = RGBColor(37, 99, 235)   # #2563EB
DARK_GRAY    = RGBColor(35, 47, 62)    # #232F3E
LIGHT_GRAY   = RGBColor(248, 249, 249) # #F8F9F9

AWS_STORAGE  = RGBColor(255, 153, 0)   # S3
AWS_DB       = RGBColor(140, 79, 255)  # DynamoDB
AWS_COMPUTE  = RGBColor(255, 181, 71)  # Lambda
AWS_NET      = RGBColor(30, 115, 190)  # CloudFront/Amplify
AWS_SECURITY = RGBColor(0, 125, 188)   # Cognito
AWS_AI       = RGBColor(254, 203, 47)  # AI/ML
AWS_MONITOR  = RGBColor(0, 153, 153)   # CloudWatch
AWS_GRAY     = RGBColor(95, 107, 109)  # connectors

# ============ Helpers ============
def add_watermark(slide, text="ScopeIQ"):
    tx = slide.shapes.add_textbox(Inches(7.0), Inches(6.7), Inches(3.0), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.color.rgb = SCOPEIQ_BLUE
    p.alignment = PP_ALIGN.RIGHT

def add_tiny_ai_icon(slide):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.25), Inches(6.8), Inches(0.35), Inches(0.35))
    circ.fill.solid()
    circ.fill.fore_color.rgb = SCOPEIQ_BLUE
    circ.fill.transparency = 0.2
    circ.line.color.rgb = SCOPEIQ_BLUE
    tx = slide.shapes.add_textbox(Inches(0.25), Inches(6.8), Inches(0.35), Inches(0.35))
    tf = tx.text_frame
    tf.text = "AI"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = DARK_GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bullets_slide(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    body = s.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i else body.paragraphs[0]
        p.text = b
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
    add_watermark(s); add_tiny_ai_icon(s)
    return s

def add_table_slide(prs, title, headers, rows):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    row_count = 1 + len(rows)
    col_count = len(headers)
    tbl = s.shapes.add_table(row_count, col_count, Inches(0.5), Inches(1.4), Inches(9.0), Inches(1.2 + 0.35*row_count)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    for i, r in enumerate(rows, start=1):
        for j, val in enumerate(r):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    add_watermark(s); add_tiny_ai_icon(s)
    return s

def add_chart_png(path, kind, x, y, title):
    plt.figure(figsize=(4,3))
    if kind == "line":
        plt.plot(x, y, marker="o")
    elif kind == "bar":
        plt.bar(x, y)
    elif kind == "area":
        plt.fill_between(x, y, alpha=0.3)
        plt.plot(x, y, marker="o")
    plt.title(title)
    plt.grid(True if kind != "bar" else False)
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()

def add_images_row(prs, title, image_paths, top_in=1.5, heights_in=2.5):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    left = Inches(0.5)
    for img in image_paths:
        s.shapes.add_picture(img, left, Inches(top_in), height=Inches(heights_in))
        left = Inches(left.inches + 3.3)
    add_watermark(s); add_tiny_ai_icon(s)
    return s

def add_service_box(slide, label, x, y, w, h, color, opacity=0.15, text_color=DARK_GRAY, bold=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.fill.transparency = opacity
    rect.line.color.rgb = color
    tx = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(h-0.2))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = label; r.font.size = Pt(12); r.font.bold = bold; r.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return rect

def add_connector(slide, x1, y1, x2, y2):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = AWS_GRAY
    ln.line.width = Pt(1.5)
    return ln

# ============ Build deck ============
prs = Presentation()

# 1) Title
title = prs.slides.add_slide(prs.slide_layouts[0])
title.shapes.title.text = "ScopeIQ — Go‑to‑Market Roadmap & Architecture"
title.placeholders[1].text = "AI‑powered insights for smarter builds • 2025"
add_watermark(title); add_tiny_ai_icon(title)

# 2) Vision & Opportunity
add_bullets_slide(prs, "Vision & Opportunity", [
    "Empower AEC & facilities teams to interpret plans/specs instantly with AI",
    "$20B+ construction tech market; pain = time wasted, errors, collaboration gaps",
    "ScopeIQ edge: AI querying, continuous ML, secure multi‑tenant SaaS"
])

# 3) Target Segments
add_bullets_slide(prs, "Target Segments", [
    "Architecture & Engineering — instant plan/spec queries",
    "Construction PMs — faster answers, fewer RFIs",
    "Facility Owners/Managers — quick retrieval from archives"
])

# 4) Value Proposition
add_bullets_slide(prs, "Value Proposition", [
    "For Users: Save time • Reduce errors • Collaborate better",
    "For Firms: Lower costs • Win more bids • Future‑proof workflows"
])

# 5) Product & Pricing
headers = ["Tier", "Price", "Who", "Key Features"]
rows = [
    ["Free", "$0", "Solo / trial", "3 uploads/mo, basic AI query"],
    ["Starter", "$49 / user", "Small teams", "Unlimited uploads, standard AI"],
    ["Pro", "$99–149 / user", "Growing teams / API", "Multi‑user, faster AI, API access"],
    ["Enterprise", "Custom", "Large orgs / SSO", "Unlimited, custom AI training, SSO, SLA"],
]
add_table_slide(prs, "Product & Pricing", headers, rows)

# 6) GTM Roadmap
add_bullets_slide(prs, "GTM Roadmap (Overview)", [
    "Pre‑Launch (M‑3→M0): pricing, onboarding, beta recruitment, marketing kit",
    "Launch (M0→M3): webinars, PR, LinkedIn ads, outreach, content",
    "Growth (M3→M12): integrations, events, referrals, enterprise sales"
])

# 7–9) KPI charts
months = ["M1","M2","M3","M4","M5","M6"]
mrr = [5, 10, 18, 30, 45, 65]   # $k
conv = [5, 6, 7, 8, 9, 10]      # %
users = [100, 250, 500, 800, 1200, 2000]

os.makedirs("charts", exist_ok=True)
mrr_img = "charts/kpi_mrr.png"
conv_img = "charts/kpi_conv.png"
users_img = "charts/kpi_users.png"
add_chart_png(mrr_img, "line", months, mrr, "MRR Growth ($k)")
add_chart_png(conv_img, "bar", months, conv, "Free → Paid Conversion (%)")
add_chart_png(users_img, "area", months, users, "User Growth")
add_images_row(prs, "Key Performance Indicators", [mrr_img, conv_img, users_img])

# 10) Launch Plan
add_bullets_slide(prs, "Launch Plan (0–3 Months)", [
    "Weekly live demos/webinars (record & publish)",
    "PR to construction tech outlets",
    "LinkedIn ads & targeted outreach (≥500 firms)",
    "2 articles/month — AI for plans/specs, case studies",
    "KPIs: 500 free sign‑ups; 5–8% conversion; CAC < 30% of Year‑1 LTV"
])

# 11) Risks & Mitigation
add_bullets_slide(prs, "Risks & Mitigation", [
    "Slow adoption → seed case studies and ROI proof",
    "Competitive AI → proprietary training + strategic integrations",
    "Integration complexity → prioritize Autodesk/Procore first"
])

# 12) Next Steps
add_bullets_slide(prs, "Next Steps", [
    "Finalize beta recruitment & early access",
    "Publish launch PR & schedule webinars",
    "Implement referral & affiliate programs",
    "Prepare 3 early customer case studies"
])

# 13) High‑Level Architecture (shapes)
s_hl = prs.slides.add_slide(prs.slide_layouts[5])
s_hl.shapes.title.text = "High‑Level Architecture"
add_service_box(s_hl, "Client Layer\nReact Web App • PWA • Query Scoping", 0.5, 1.2, 9.0, 0.9, LIGHT_GRAY, opacity=1.0)
add_service_box(s_hl, "Amplify / CloudFront", 0.5, 2.4, 2.7, 1.0, AWS_NET)
add_service_box(s_hl, "Cognito", 3.35, 2.4, 2.2, 1.0, AWS_SECURITY)
add_service_box(s_hl, "AppSync", 5.8, 2.4, 2.2, 1.0, SCOPEIQ_BLUE, opacity=0.25, text_color=RGBColor(255,255,255))
add_service_box(s_hl, "Lambda", 8.2, 2.4, 1.3, 1.0, AWS_COMPUTE)
add_service_box(s_hl, "Amazon S3\n(Documents & Thumbnails)", 0.5, 3.7, 3.2, 1.0, AWS_STORAGE)
add_service_box(s_hl, "DynamoDB\n(Multi‑Tenant DB)", 3.9, 3.7, 3.0, 1.0, AWS_DB)
add_service_box(s_hl, "AI & Search\nOpenAI • Pinecone • PDF.js", 7.1, 3.7, 2.4, 1.0, AWS_AI)
add_service_box(s_hl, "CloudWatch\n(Monitoring & Logs)", 0.5, 4.9, 2.7, 0.9, AWS_MONITOR)
add_connector(s_hl, 1.9, 2.4, 1.9, 1.2)
add_connector(s_hl, 4.45, 2.4, 4.45, 1.2)
add_connector(s_hl, 6.9, 2.4, 6.9, 1.2)
add_connector(s_hl, 6.9, 3.4, 6.9, 2.4)
add_connector(s_hl, 2.1, 3.7, 2.1, 3.4)
add_connector(s_hl, 5.4, 3.7, 5.4, 3.4)
add_connector(s_hl, 8.3, 3.4, 8.3, 2.4)
add_connector(s_hl, 8.3, 4.7, 8.3, 4.3)
tx = s_hl.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9.0), Inches(0.7))
tf = tx.text_frame; tf.text = ("Amplify/CloudFront deliver the app; Cognito secures access; AppSync provides APIs; "
                               "Lambda handles processing; S3/DynamoDB store content; AI/Search (OpenAI, Pinecone, PDF.js) "
                               "powers insights; CloudWatch monitors.")
tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].font.color.rgb = DARK_GRAY
add_watermark(s_hl); add_tiny_ai_icon(s_hl)

# 14) AI & Data Flow (shapes + AI highlight)
s_flow = prs.slides.add_slide(prs.slide_layouts[5])
s_flow.shapes.title.text = "AI & Data Flow"
def box(text, x, y, w, h, col, op=0.15, white=False):
    return add_service_box(s_flow, text, x, y, w, h, col, opacity=op, text_color=RGBColor(255,255,255) if white else DARK_GRAY)
box("1) Upload → S3", 0.5, 1.6, 2.5, 0.9, AWS_STORAGE)
box("2) Metadata → DynamoDB", 3.2, 1.6, 2.5, 0.9, AWS_DB)
box("3) Processing → Lambda", 5.9, 1.6, 2.5, 0.9, AWS_COMPUTE)
add_service_box(s_flow, "4) AI Analysis → Textract • Comprehend • GPT‑4 • PDF.js", 0.5, 2.8, 8.0, 0.9, AWS_AI, opacity=0.10)
box("5) Embedding/Index → Pinecone", 0.5, 4.0, 2.9, 0.9, SCOPEIQ_BLUE, op=0.25, white=True)
box("6) Search API → AppSync", 3.6, 4.0, 2.9, 0.9, SCOPEIQ_BLUE, op=0.25, white=True)
box("7) Monitoring → CloudWatch", 6.7, 4.0, 2.0, 0.9, AWS_MONITOR)
add_connector(s_flow, 1.75, 2.5, 4.45, 2.5)
add_connector(s_flow, 4.45, 2.5, 7.15, 2.5)
add_connector(s_flow, 7.15, 2.5, 7.15, 2.8)
add_connector(s_flow, 2.0, 4.9, 4.0, 4.9)
tx2 = s_flow.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9.0), Inches(0.7))
tf2 = tx2.text_frame; tf2.text = ("S3 stores files; DynamoDB tracks metadata; Lambda orchestrates analysis; "
                                  "embeddings go to Pinecone; AppSync exposes semantic search; CloudWatch monitors.")
tf2.paragraphs[0].font.size = Pt(12); tf2.paragraphs[0].font.color.rgb = DARK_GRAY
add_watermark(s_flow); add_tiny_ai_icon(s_flow)

# 15) Appendix — Simplified Technical Architecture (shapes)
s_appx = prs.slides.add_slide(prs.slide_layouts[5])
s_appx.shapes.title.text = "Appendix — Technical Architecture (Simplified)"
add_service_box(s_appx, "Client Layer\n(React, PWA, Query Scoping)", 0.5, 1.3, 3.0, 0.9, LIGHT_GRAY, opacity=1.0)
add_service_box(s_appx, "Hosting/CDN\n(Amplify, CloudFront)", 3.7, 1.3, 2.5, 0.9, AWS_NET)
add_service_box(s_appx, "Auth\n(Cognito, JWT, RBAC)", 6.4, 1.3, 2.5, 0.9, AWS_SECURITY)
add_service_box(s_appx, "API Layer\n(AppSync)", 0.5, 2.5, 2.5, 0.9, SCOPEIQ_BLUE, opacity=0.25, text_color=RGBColor(255,255,255))
add_service_box(s_appx, "Processing\n(Lambda)", 3.2, 2.5, 2.0, 0.9, AWS_COMPUTE)
add_service_box(s_appx, "Storage\n(S3 Docs/Thumbs)", 5.5, 2.5, 2.0, 0.9, AWS_STORAGE)
add_service_box(s_appx, "Database\n(DynamoDB)", 7.7, 2.5, 1.2, 0.9, AWS_DB)
add_service_box(s_appx, "AI & Search\n(GPT‑4, Textract, Comprehend, Pinecone, PDF.js)", 0.5, 3.7, 5.6, 0.9, AWS_AI, opacity=0.10)
add_service_box(s_appx, "Monitoring\n(CloudWatch)", 6.3, 3.7, 2.6, 0.9, AWS_MONITOR)
add_watermark(s_appx); add_tiny_ai_icon(s_appx)

# Save deck
out = "ScopeIQ_GTM_Architecture_FULL.pptx"
prs.save(out)
print(f"Saved {out}")
